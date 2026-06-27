import os
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation


def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF using PyMuPDF."""
    doc = fitz.open(file_path)
    text = ""
    for page in doc:
        text += page.get_text()
    doc.close()
    return text.strip()


def extract_text_from_docx(file_path: str) -> str:
    """Extract text from DOCX file."""
    doc = Document(file_path)
    paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
    return "\n".join(paragraphs).strip()


def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PPTX file."""
    prs = Presentation(file_path)
    text = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        if slide_text:
            text.append(f"[Slide {slide_num}]\n" + "\n".join(slide_text))
    return "\n\n".join(text).strip()


def extract_text_from_code(file_path: str) -> str:
    """Extract text from code files."""
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read().strip()


def extract_text(file_path: str) -> dict:
    """
    Main extractor — auto-detects file type and extracts text.
    Returns a dict with extracted text, file type, and metadata.
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    ext = os.path.splitext(file_path)[1].lower()
    file_name = os.path.basename(file_path)
    file_size = os.path.getsize(file_path)

    extractors = {
        ".pdf":  extract_text_from_pdf,
        ".docx": extract_text_from_docx,
        ".pptx": extract_text_from_pptx,
        ".py":   extract_text_from_code,
        ".js":   extract_text_from_code,
        ".java": extract_text_from_code,
        ".cpp":  extract_text_from_code,
        ".c":    extract_text_from_code,
        ".txt":  extract_text_from_code,
    }

    if ext not in extractors:
        raise ValueError(f"Unsupported file type: {ext}. Supported: {list(extractors.keys())}")

    raw_text = extractors[ext](file_path)

    return {
        "file_name": file_name,
        "file_path": file_path,
        "file_type": ext,
        "file_size_kb": round(file_size / 1024, 2),
        "word_count": len(raw_text.split()),
        "char_count": len(raw_text),
        "raw_text": raw_text,
    }
